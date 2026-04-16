"""JSON to PPTX conversion — local file I/O version.

Ported from: ppt-addin/backend/services/make_pptx/slide_gen/slide_builder.py
Changes: Azure Blob → local file system, template bundled in package.
"""
import os
from io import BytesIO
from pathlib import Path
from typing import Optional

from PIL import Image
from pptx import Presentation

from slidegenie.utils.constants import PowerPointConfig
from slidegenie.shapes.object_function import add_shapes_to_slide
from slidegenie.utils.common import replace_texts_by_shape_name


def _get_template_path() -> str:
    """Get path to the bundled PPTX template."""
    return str(Path(__file__).parent.parent / "templates" / "template.pptx")


def json_to_pptx(
    make_type: str,
    logger,
    json_data: dict,
    src_image: Image.Image | None = None,
    output_path: str = "output.pptx",
) -> str:
    """Convert OCR JSON data to a PPTX file.

    Args:
        make_type: Slide type (graphic, flow, matrix)
        logger: Logger instance
        json_data: OCR JSON data with title, lead, body
        src_image: Source image to embed as thumbnail (top-right)
        output_path: Path to save the output PPTX

    Returns:
        Path to the saved PPTX file
    """
    if make_type not in {"graphic", "flow", "matrix"}:
        logger.error(f"Unsupported make_type: {make_type}")
        raise ValueError(f"Unsupported make_type: {make_type}")

    logger.info(f"title: {json_data.get('title', '')}")
    logger.info(f"lead sentence: {json_data.get('lead', '')}")

    slide_index = 0
    ppi_1k = PowerPointConfig.PPI_1K

    # Load template
    template_path = _get_template_path()
    prs = Presentation(template_path)

    try:
        prs.slides[slide_index]
    except IndexError:
        logger.error("[ERROR] slide does not exist in template.")
        raise

    # Replace title and lead text
    change_writing_contents = [
        {"title": json_data.get("title") or ""},
        {"lead": json_data.get("lead") or ""},
    ]
    replace_texts_by_shape_name(
        prs=prs,
        slide_index=slide_index,
        replacements=change_writing_contents,
    )

    # Convert coordinates from pixels to inches
    shapes_data = json_data.get("body", [])
    shapes_data = [
        {k: (v / ppi_1k if k in ["x", "y", "width", "height"] else v) for k, v in shape.items()}
        for shape in shapes_data
    ]

    # Embed source image as thumbnail in top-right corner
    image_bytesio: Optional[BytesIO] = None
    if src_image is not None:
        image_bytesio = BytesIO()
        src_image.save(image_bytesio, format="PNG")
        image_bytesio.seek(0)

        shapes_data.append({
            "image_path": image_bytesio,
            "x": PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.AI_ICON_POSITION_X_RATIO,
            "y": PowerPointConfig.AI_ICON_Y_POSITION,
            "width": PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.AI_ICON_SCALE_RATIO,
            "height": PowerPointConfig.SLIDE_HEIGHT_INCHES * PowerPointConfig.AI_ICON_SCALE_RATIO,
            "shape_type": "AI_ICON",
        })

    # Add shapes to slide
    prs = add_shapes_to_slide(shapes_data, prs, slide_index, logger=logger)

    # Release image BytesIO
    if image_bytesio is not None:
        image_bytesio.close()

    # Save locally
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    logger.info(f"PPTX saved: {output_path}")

    return output_path


def _prepare_template_slide(prs: Presentation, slide_index: int):
    """Prepare a slide to match the template's first slide.

    For slide 0 (the template itself): nothing to do, already has "title"/"lead".
    For slides 1+: remove layout auto-placeholders, then add textboxes
    with "title" and "lead" text at the exact same positions as the template.

    Template shape positions (EMU):
      "タイトル 5":                left=577146,  top=128431,  width=10915201, height=495024
      "コンテンツ プレースホルダー 2": left=514800,  top=837655,  width=11163600, height=360000
    """
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    if slide_index == 0:
        return  # Template slide already has the shapes

    slide = prs.slides[slide_index]

    # Remove auto-generated placeholders from layout
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Add title textbox matching template exactly
    txBox = slide.shapes.add_textbox(
        Emu(577146), Emu(128431), Emu(10915201), Emu(495024)
    )
    txBox.name = "タイトル 5"
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "title"
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.name = "Yu Gothic Light"
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)

    # Add lead textbox matching template exactly
    txBox = slide.shapes.add_textbox(
        Emu(514800), Emu(837655), Emu(11163600), Emu(360000)
    )
    txBox.name = "コンテンツ プレースホルダー 2"
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "lead"
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.name = "Yu Gothic Light"
    run.font.color.rgb = RGBColor(80, 80, 80)


def _render_slide(prs: Presentation, slide_index: int, logger,
                  json_data: dict, src_image: Image.Image | None = None):
    """Render JSON data onto a slide.

    All slides have "title" and "lead" text shapes (either from template
    or created by _prepare_template_slide), so replace_texts_by_shape_name
    works uniformly on all slides.
    """
    ppi_1k = PowerPointConfig.PPI_1K

    # Replace "title" and "lead" text (works on all slides uniformly)
    replace_texts_by_shape_name(
        prs=prs,
        slide_index=slide_index,
        replacements=[
            {"title": json_data.get("title") or ""},
            {"lead": json_data.get("lead") or ""},
        ],
    )

    # Convert coordinates from pixels to inches
    shapes_data = json_data.get("body", [])
    shapes_data = [
        {k: (v / ppi_1k if k in ["x", "y", "width", "height"] else v) for k, v in shape.items()}
        for shape in shapes_data
    ]

    # Embed source image as thumbnail in top-right corner
    image_bytesio: BytesIO | None = None
    if src_image is not None:
        image_bytesio = BytesIO()
        src_image.save(image_bytesio, format="PNG")
        image_bytesio.seek(0)
        shapes_data.append({
            "image_path": image_bytesio,
            "x": PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.AI_ICON_POSITION_X_RATIO,
            "y": PowerPointConfig.AI_ICON_Y_POSITION,
            "width": PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.AI_ICON_SCALE_RATIO,
            "height": PowerPointConfig.SLIDE_HEIGHT_INCHES * PowerPointConfig.AI_ICON_SCALE_RATIO,
            "shape_type": "AI_ICON",
        })

    # Add shapes to slide
    add_shapes_to_slide(shapes_data, prs, slide_index, logger=logger)

    if image_bytesio is not None:
        image_bytesio.close()


def multi_json_to_pptx(
    slides_data: list[dict],
    logger,
    output_path: str = "output.pptx",
) -> str:
    """Build a single PPTX with multiple slides from a list of JSON data.

    Strategy:
    1. Load template.pptx (preserves slide master, layout, formatting)
    2. Slide 1: use existing template slide (already has "title"/"lead")
    3. Slides 2-N: add_slide from same layout, then add "title"/"lead"
       textboxes at the exact same positions as the template
    4. replace_texts_by_shape_name + add_shapes_to_slide on ALL slides uniformly
    """
    template_path = _get_template_path()
    prs = Presentation(template_path)

    # Get layout from the template's existing slide
    template_layout = prs.slides[0].slide_layout
    logger.info(f"Using layout: \"{template_layout.name}\"")

    # Pre-create all additional slides from the same layout
    for _ in range(len(slides_data) - 1):
        prs.slides.add_slide(template_layout)

    # Prepare each slide to match the template (add "title"/"lead" textboxes)
    for i in range(len(slides_data)):
        _prepare_template_slide(prs, i)

    logger.info(f"Prepared {len(prs.slides)} slides")

    # Render content into each slide (all uniform - replace "title"/"lead" + add body)
    for i, slide_info in enumerate(slides_data):
        json_data = slide_info["json_data"]
        src_image = slide_info.get("src_image")
        logger.info(f"[{i + 1}/{len(slides_data)}] Rendering: {json_data.get('title', '')[:50]}")
        _render_slide(prs, i, logger, json_data, src_image)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    logger.info(f"Multi-slide PPTX saved: {output_path} ({len(prs.slides)} slides)")
    return output_path
