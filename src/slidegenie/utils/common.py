"""Utility functions for slidegenie.

Ported from: ppt-addin/backend/services/common_function.py
"""
import os
import re
import base64
from io import BytesIO
from typing import Optional

from pptx.presentation import Presentation
from jinja2 import Environment, FileSystemLoader
from PIL import Image, ImageOps

from slidegenie.utils.constants import PowerPointConfig


def decode_base64_image(image_base64: str) -> Image.Image:
    """Decode Base64-encoded image string to PIL Image."""
    try:
        if "," in image_base64:
            image_base64 = image_base64.split(",")[1]
        image_data = base64.b64decode(image_base64)
        image = Image.open(BytesIO(image_data))
        return image
    except Exception as e:
        raise ValueError(f"Failed to decode image: {str(e)}")


def replace_texts_by_shape_name(
    prs: Presentation,
    slide_index: int,
    replacements: list[dict[str, str]],
):
    """Replace text in slide shapes matching keywords."""
    slide = prs.slides[slide_index]

    for rep in replacements:
        for keyword, new_text in rep.items():
            found = False
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text.strip()
                        if original_text == keyword:
                            run.text = new_text
                            found = True
                            break
                    if found:
                        break
                if found:
                    break


def load_prompt(template_name: str, **kwargs) -> str:
    """Load and render a Jinja2 prompt template.

    Searches for templates in the prompts/ directory within the slidegenie package.
    """
    base_dir = os.path.dirname(os.path.dirname(__file__))
    template_dirs = [
        os.path.join(base_dir, "prompts", "tone-manner"),
        os.path.join(base_dir, "prompts", "addin"),
    ]

    env = Environment(
        loader=FileSystemLoader(template_dirs),
        trim_blocks=True,
        lstrip_blocks=True,
    )

    fname = f"{template_name}.j2"
    template = env.get_template(fname)
    return template.render(**kwargs)


def resolve_image_prompt_template(
    *,
    english_frag: bool,
    addin_en_template: str,
    addin_ja_template: str,
) -> str:
    """Resolve prompt template name by language."""
    return addin_en_template if english_frag else addin_ja_template


def normalize_to_fullhd(img: Image.Image) -> Image.Image:
    """Normalize image to Full HD (1920x1080)."""
    img = ImageOps.exif_transpose(img)

    if img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGB")

    W, H = PowerPointConfig.FULLHD_W, PowerPointConfig.FULLHD_H

    img = ImageOps.contain(img, (W, H), method=Image.LANCZOS)

    canvas = Image.new("RGB", (W, H), (255, 255, 255))
    x = (W - img.size[0]) // 2
    y = (H - img.size[1]) // 2
    canvas.paste(img, (x, y))
    return canvas


def is_english(text: str) -> bool:
    """Detect if text is primarily English."""
    jp_pattern = re.compile(r'[ぁ-んァ-ン一-龥]')
    if jp_pattern.search(text):
        return False
    alpha_count = sum(1 for c in text if c.isalpha() and c.isascii())
    total = len(text.strip())
    if total == 0:
        return False
    return (alpha_count / total) > 0.3
