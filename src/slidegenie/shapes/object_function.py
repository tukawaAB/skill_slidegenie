import unicodedata
from typing import Any, Dict, List, Optional, Union
from logging import Logger

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from slidegenie.utils.constants import PowerPointConfig


# Shape properties mapping
SHAPE_PROPERTIES = {
    "required": PowerPointConfig.REQUIRED_PROPERTIES,
    "text": {
        "text": None,
        "font_size": PowerPointConfig.DEFAULT_FONT_SIZE,
        "font_name": PowerPointConfig.DEFAULT_FONT_NAME,
        "font_color": PowerPointConfig.DEFAULT_FONT_COLOR,
        "text_align": PowerPointConfig.DEFAULT_TEXT_ALIGNMENT,
    },
}


# ============================================================================
# CHARACTER AND TEXT ANALYSIS
# ============================================================================


def get_standard_char_width_ratio() -> float:
    """Get standard character width ratio for text calculations.

    Centralized function to avoid using CJK_DEFAULT_CHAR_WIDTH_RATIO directly.
    """
    return PowerPointConfig.CJK_DEFAULT_CHAR_WIDTH_RATIO


def get_shape_text_area_ratio(shape_type_name: str) -> float:
    """
    Get text area ratio for different shape types.
    Args:
        shape_type_name: Shape type name (e.g., "OVAL", "PENTAGON", "RECTANGLE")

    Returns:
        Ratio of usable text area (0.0 to 1.0)
    """
    return PowerPointConfig.SHAPE_TEXT_AREA_RATIOS.get(
        shape_type_name.upper(), PowerPointConfig.SHAPE_TEXT_AREA_RATIOS["DEFAULT"]
    )


def calculate_font_metrics(font_size_pt: float) -> Dict[str, float]:
    """Calculate font metrics for text dimension calculations.

    Simplified version that uses standard char width ratio for all text.
    """
    font_size_inches = font_size_pt / PowerPointConfig.POINTS_PER_INCH
    char_width_ratio = get_standard_char_width_ratio()

    return {
        "font_size_inches": font_size_inches,
        "char_width_inches": font_size_inches * char_width_ratio,
        "line_height": font_size_inches * PowerPointConfig.LINE_HEIGHT_RATIO,
    }


def _is_full_width_character(char: str) -> bool:
    """Check if character is full-width using Unicode classification."""
    return unicodedata.east_asian_width(char) in ("F", "W")


def _is_half_width_character(char: str) -> bool:
    """Check if character is half-width using Unicode classification."""
    width = unicodedata.east_asian_width(char)
    return width in ("H", "Na", "N")


def detect_text_language(text: str) -> str:
    """Detect if text contains CJK characters using Unicode classification."""
    if not text:
        return "latin"

    for char in text:
        if _is_full_width_character(char):
            return "cjk"

        # Check specific CJK ranges
        char_code = ord(char)
        if (
            0x4E00 <= char_code <= 0x9FFF  # CJK Unified Ideographs
            or 0x3400 <= char_code <= 0x4DBF  # CJK Extension A
            or 0x3040 <= char_code <= 0x309F  # Hiragana
            or 0x30A0 <= char_code <= 0x30FF
        ):  # Katakana
            return "cjk"

    return "latin"


def get_character_weight(char: str) -> float:
    """Get character weight for text calculations using Unicode classification."""
    if _is_full_width_character(char):
        return PowerPointConfig.FULL_WIDTH_CHAR_WIDTH_RATIO
    elif char.isspace():
        return PowerPointConfig.SPACE_CHAR_WIDTH_RATIO
    elif _is_half_width_character(char):
        return (
            PowerPointConfig.CJK_CHAR_WIDTH_RATIO
            if detect_text_language(char) == "cjk"
            else PowerPointConfig.HALF_WIDTH_CHAR_WIDTH_RATIO
        )
    else:
        return PowerPointConfig.CJK_CHAR_WIDTH_RATIO


def calculate_text_weight(text: str) -> float:
    """Calculate total weighted text length for all text calculations."""
    if not text:
        return PowerPointConfig.OBJECT_FUNCTION_CONFIG["MIN_CHAR_WEIGHT"]
    return sum(get_character_weight(char) for char in text)


# Backward compatibility alias
calculate_weighted_text_length = calculate_text_weight

# ============================================================================
# TEXT APPLICATION TO SHAPES
# ============================================================================

def apply_text_to_shape(shape, text: str, shape_data: Dict[str, Any], logger: Logger = None) -> None:
    """Apply text to shape with 0.1 inch padding and proper formatting."""
    try:
        if not hasattr(shape, "text_frame"):
            if logger:
                logger.warning("Shape does not support text")
            return

        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        # Set 0.1 inch margins for padding (as per requirement)
        text_frame.margin_left = Inches(PowerPointConfig.TEXT_PADDING_INCHES)
        text_frame.margin_right = Inches(PowerPointConfig.TEXT_PADDING_INCHES)
        text_frame.margin_top = Inches(PowerPointConfig.TEXT_PADDING_INCHES)
        text_frame.margin_bottom = Inches(PowerPointConfig.TEXT_PADDING_INCHES)

        # Disable auto-size to prevent PowerPoint from overriding settings
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

        # Get formatting properties once
        font_size = shape_data.get("font_size", PowerPointConfig.DEFAULT_FONT_SIZE)
        font_size_pt = Pt(font_size)
        font_name = shape_data.get("font_name", PowerPointConfig.DEFAULT_FONT_NAME)
        font_color_rgb = (
            RGBColor(*shape_data["font_color"])
            if "font_color" in shape_data and shape_data["font_color"]
            else None
        )
        font_style = shape_data.get("font_style")

        alignment_map = {
            "center": PP_ALIGN.CENTER,
            "left": PP_ALIGN.LEFT,
            "right": PP_ALIGN.RIGHT,
        }
        lines = text.split("\n")        
        text_align_str = shape_data.get("text_align", "center")
        alignment = alignment_map.get(text_align_str, PP_ALIGN.CENTER)

        # Handle text line by line to ensure consistent formatting
        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = line
            p.alignment = alignment

            # Ensure run exists and apply formatting
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = font_size_pt
            # text-level shadowを除去: rPr に空の effectLst を追加してテーマ継承をブロック
            rpr = run._r.get_or_add_rPr()
            for existing in rpr.findall(qn('a:effectLst')):
                rpr.remove(existing)
            rpr.append(OxmlElement('a:effectLst'))
            if font_name:
                run.font.name = str(font_name)
            if font_color_rgb:
                run.font.color.rgb = font_color_rgb
            # Apply font style
            if font_style:
                if "bold" in font_style:
                    run.font.bold = True
                if "italic" in font_style:
                    run.font.italic = True

        cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
        if logger:
            logger.debug(f"Applied text to shape: '{text[:cfg['TEXT_PREVIEW_LENGTH']]}...' Font: {font_size}pt")

    except Exception as e:
        if logger:
            logger.error(f"Failed to apply text: {e}")
        # Simple fallback
        try:
            shape.text_frame.text = text
        except Exception:
            if logger:
                logger.error("Fallback text application failed")


# ============================================================================
# TEXT TRUNCATION LOGIC
# ============================================================================


def truncate_text_with_ellipsis(
    text: str,
    font_size_pt: float,
    shape_width_inches: float,
    shape_height_inches: float,
    shape_type_name: str = "RECTANGLE",
    is_truncate: bool = False,
) -> str:
    """Truncate text with ellipsis to fit within shape boundaries."""
    if not text or not is_truncate:
        return text

    # Auto-detect language
    language = detect_text_language(text)

    # Calculate capacity
    capacity = calculate_text_capacity(
        font_size_pt, shape_width_inches, shape_height_inches, shape_type_name
    )

    # Early return if capacity too small
    if capacity <= PowerPointConfig.ELLIPSIS_CHAR_COUNT:
        return PowerPointConfig.ELLIPSIS_TEXT

    # Check if text fits
    text_weight = calculate_text_weight(text)
    if text_weight <= capacity:
        return text

    # Truncate text - reserve space for ellipsis
    target_capacity = capacity - PowerPointConfig.ELLIPSIS_CHAR_COUNT

    return (
        _truncate_cjk_text(text, target_capacity)
        if language == "cjk"
        else _truncate_latin_text(text, target_capacity)
    )


def _truncate_cjk_text(text: str, target_capacity: float) -> str:
    """Truncate CJK text by character with accurate weights."""
    cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
    current_length = cfg["MIN_CHAR_WEIGHT"]
    truncate_index = cfg["MIN_TRUNCATE_INDEX"]

    for i, char in enumerate(text):
        char_weight = get_character_weight(char)
        if current_length + char_weight > target_capacity:
            break
        current_length += char_weight
        truncate_index = i + 1

    if truncate_index == cfg["MIN_TRUNCATE_INDEX"]:
        return PowerPointConfig.ELLIPSIS_TEXT
    return text[:truncate_index] + PowerPointConfig.ELLIPSIS_TEXT


def _truncate_latin_text(text: str, target_capacity: float) -> str:
    """Truncate Latin text at word boundaries with accurate weights."""
    cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
    words = text.split()
    text_parts = []
    current_length = cfg["MIN_CHAR_WEIGHT"]

    for word in words:
        # Calculate space + word weight
        space_weight = PowerPointConfig.SPACE_CHAR_WIDTH_RATIO if text_parts else 0
        word_weight = calculate_text_weight(word)
        total_length = current_length + space_weight + word_weight

        if total_length > target_capacity:
            break

        text_parts.append(word)
        current_length = total_length

    # If no words fit, truncate character by character
    if not text_parts:
        return _truncate_cjk_text(text, target_capacity)

    return " ".join(text_parts) + PowerPointConfig.ELLIPSIS_TEXT


# ============================================================================
# TEXT DIMENSION CALCULATIONS WITH SHAPE RATIOS
# ============================================================================


def _calculate_latin_text_lines(text: str, chars_per_line: int) -> int:
    """Calculate lines needed for Latin text with word wrapping."""
    cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
    words = text.split()
    if not words:
        return cfg["MIN_LINES_NEEDED"]

    lines_needed = cfg["MIN_LINES_NEEDED"]
    current_line_length = 0

    for word in words:
        word_length = len(word)
        # Check if adding word would exceed line length
        if (
            current_line_length > 0
            and (current_line_length + cfg["SPACE_CHAR_WEIGHT_OFFSET"] + word_length) > chars_per_line
        ):
            lines_needed += 1
            current_line_length = word_length
        else:
            if current_line_length > 0:
                current_line_length += cfg["SPACE_CHAR_WEIGHT_OFFSET"]  # space
            current_line_length += word_length

    return lines_needed


def estimate_text_dimensions(
    text: str, font_size_pt: float, shape_width_inches: float, shape_type_name: str
) -> Dict[str, Union[float, int]]:
    """Estimate text dimensions with 0.05 inch padding, word wrapping and shape ratio."""
    cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
    if not text:
        return {
            "width_inches": cfg["MIN_TEXT_WIDTH"], 
            "height_inches": cfg["MIN_TEXT_HEIGHT"], 
            "lines": cfg["MIN_TEXT_LINES"]
        }

    # Get font metrics
    font_metrics = calculate_font_metrics(font_size_pt)

    # Apply 0.05 inch padding on both sides (total 0.1 inch reduction)
    available_width = shape_width_inches - (cfg["PADDING_REDUCTION_FACTOR"] * PowerPointConfig.TEXT_PADDING_INCHES)
    if available_width <= 0:
        return {
            "width_inches": cfg["MIN_TEXT_WIDTH"],
            "height_inches": font_metrics["line_height"],
            "lines": cfg["MAX_TEXT_LINES"],
        }

    # Apply shape-specific text area ratio
    shape_ratio = get_shape_text_area_ratio(shape_type_name)
    effective_width = available_width * shape_ratio

    # Calculate characters per line using effective width
    chars_per_line = max(cfg["MIN_CHARS_PER_LINE"], int(effective_width / font_metrics["char_width_inches"]))

    # Calculate lines needed based on language
    language = detect_text_language(text)
    if language == "cjk":
        # CJK text can break anywhere
        text_weight = calculate_text_weight(text)
        lines_needed = max(cfg["MIN_LINES_NEEDED"], (text_weight + chars_per_line - 1) // chars_per_line)
    else:
        # Latin text needs word boundary wrapping
        lines_needed = _calculate_latin_text_lines(text, chars_per_line)

    # Calculate final dimensions
    text_height = lines_needed * font_metrics["line_height"]

    # Estimate text width
    if lines_needed == cfg["MAX_TEXT_LINES"]:
        text_weight = calculate_text_weight(text)
        text_width = min(
            text_weight * font_metrics["char_width_inches"], effective_width
        )
    else:
        text_width = effective_width  # Multi-line uses full effective width

    return {
        "width_inches": text_width,
        "height_inches": text_height,
        "lines": lines_needed,
    }


def calculate_text_capacity(
    font_size_pt: float,
    shape_width_inches: float,
    shape_height_inches: float,
    shape_type_name: str = None,
) -> float:
    """Calculate text capacity with 0.05 inch padding and shape-specific area ratio."""
    cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
    if shape_type_name is None:
        shape_type_name = cfg["DEFAULT_SHAPE_TYPE"]
        
    # Get font metrics
    font_metrics = calculate_font_metrics(font_size_pt)

    # Apply 0.05 inch padding on all sides
    available_width = shape_width_inches - (cfg["PADDING_REDUCTION_FACTOR"] * PowerPointConfig.TEXT_PADDING_INCHES)
    available_height = shape_height_inches - (cfg["PADDING_REDUCTION_FACTOR"] * PowerPointConfig.TEXT_PADDING_INCHES)

    if available_width <= 0 or available_height <= 0:
        return cfg["MIN_CHAR_WEIGHT"]

    # Apply shape-specific text area ratio
    shape_ratio = get_shape_text_area_ratio(shape_type_name)
    effective_width = available_width * shape_ratio
    effective_height = available_height * shape_ratio

    # Calculate capacity using effective dimensions
    chars_per_line = effective_width / font_metrics["char_width_inches"]
    lines_available = effective_height / font_metrics["line_height"]

    return chars_per_line * lines_available


# ============================================================================
# TEXT FITTING AND OPTIMIZATION
# ============================================================================


def calculate_optimal_text_and_shape(
    text: str,
    original_width_inches: float,
    original_height_inches: float,
    shape_type_name: str = None,
    initial_font_size: float = PowerPointConfig.DEFAULT_FONT_SIZE,
) -> Dict[str, Union[float, bool, str, Dict[str, Union[float, int]]]]:
    """
    Calculate optimal font size with 0.1 inch padding and shape-specific ratios.

    Key requirements:
    1. Text must fit inside shape with 0.1 inch padding
    2. If text doesn't fit at 8pt, truncate with "..."
    3. Consider full/half/latin character weights
    4. Use shape-specific text area ratios

    Args:
        text: Text content to fit
        original_width_inches: Original shape width in inches
        original_height_inches: Original shape height in inches
        shape_type_name: Shape type name for ratio calculation
        initial_font_size: Starting font size to try

    Returns:
        Dictionary containing optimized parameters and metadata
    """
    cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
    if shape_type_name is None:
        shape_type_name = cfg["DEFAULT_SHAPE_TYPE"]
        
    if not text:
        return {
            "font_size": initial_font_size,
            "shape_width": original_width_inches,
            "shape_height": original_height_inches,
            "fits": cfg["TEXT_FITS_DEFAULT"],
            "resized": cfg["RESIZED_DEFAULT"],
            "dimensions": {
                "width_inches": cfg["MIN_TEXT_WIDTH"], 
                "height_inches": cfg["MIN_TEXT_HEIGHT"], 
                "lines": cfg["MIN_TEXT_LINES"]
            },
            "auto_calculation_disabled": cfg["AUTO_CALCULATION_DISABLED"],
        }

    language = detect_text_language(text)
    text_weight = calculate_text_weight(text)
    current_font_size = initial_font_size

    # Try font sizes from initial down to minimum
    while current_font_size >= PowerPointConfig.MIN_FONT_SIZE:
        # Calculate text capacity
        capacity = calculate_text_capacity(
            current_font_size,
            original_width_inches,
            original_height_inches,
            shape_type_name,
        )

        if text_weight <= capacity:
            # Text fits! Calculate actual dimensions for return
            dimensions = estimate_text_dimensions(
                text, current_font_size, original_width_inches, shape_type_name
            )
            return {
                "font_size": current_font_size,
                "shape_width": original_width_inches,
                "shape_height": original_height_inches,
                "fits": cfg["TEXT_FITS_DEFAULT"],
                "resized": cfg["RESIZED_DEFAULT"],
                "dimensions": dimensions,
                "capacity": capacity,
                "auto_calculation_disabled": cfg["AUTO_CALCULATION_DISABLED"],
            }

        # Text doesn't fit, decrease font size by one step
        current_font_size -= PowerPointConfig.FONT_SIZE_STEP


    truncated_text = truncate_text_with_ellipsis(
        text,
        PowerPointConfig.MIN_FONT_SIZE,
        original_width_inches,
        original_height_inches,
        shape_type_name,
    )

    # Calculate dimensions of truncated text
    truncated_dimensions = estimate_text_dimensions(
        truncated_text,
        PowerPointConfig.MIN_FONT_SIZE,
        original_width_inches,
        shape_type_name,
    )

    # Get final capacity for debugging
    final_capacity = calculate_text_capacity(
        PowerPointConfig.MIN_FONT_SIZE,
        original_width_inches,
        original_height_inches,
        shape_type_name,
    )

    return {
        "font_size": PowerPointConfig.MIN_FONT_SIZE,
        "shape_width": original_width_inches,
        "shape_height": original_height_inches,
        "fits": cfg["TEXT_FITS_DEFAULT"],  # Always fits since we truncated
        "resized": cfg["RESIZED_DEFAULT"],
        "dimensions": truncated_dimensions,
        "text_truncated": True,
        "truncated_text": truncated_text,
        "capacity": final_capacity,
    }


def validate_shape_data(shape_data: Dict[str, Any], logger: Logger = None, max_font_size: Optional[int] = None) -> Optional[Dict[str, Any]]:
    """
    Validate shape data format and values, convert shape_type from string to enum.

    Args:
        shape_data: Shape data dictionary to validate and process
        logger: Logger instance for error logging
        max_font_size: Optional custom max font size (for PDF2PPTX: 72)

    Returns:
        Processed shape_data with converted shape_type, or None if validation fails
    """
    # Fix negative dimensions
    if shape_data.get("width", 0) < 0:
        shape_data["width"] = abs(shape_data["width"])
    
    if shape_data.get("height", 0) < 0:
        shape_data["height"] = abs(shape_data["height"])

    try:
        # Validate required properties exist
        if not _validate_required_properties(shape_data, logger):
            if logger:
                logger.error(f"[validate] Failed: required_properties, text='{str(shape_data.get('text', ''))[:30]}...'")
            return None

        # Validate and convert shape_type
        if not _validate_and_convert_shape_type(shape_data, logger):
            if logger:
                logger.error(f"[validate] Failed: shape_type, text='{str(shape_data.get('text', ''))[:30]}...'")
            return None

        # Validate numeric properties in batch
        if not _validate_numeric_properties(shape_data, logger):
            if logger:
                logger.error(f"[validate] Failed: numeric_properties, text='{str(shape_data.get('text', ''))[:30]}...'")
            return None

        # Validate boundaries
        if not _validate_shape_boundaries(shape_data, logger):
            if logger:
                logger.error(f"[validate] Failed: boundaries, x={shape_data.get('x')}, y={shape_data.get('y')}, w={shape_data.get('width')}, h={shape_data.get('height')}, text='{str(shape_data.get('text', ''))[:30]}...'")
            return None

        # Validate optional properties
        if not _validate_optional_properties(shape_data, logger, max_font_size=max_font_size):
            if logger:
                logger.error(f"[validate] Failed: optional_properties, text='{str(shape_data.get('text', ''))[:30]}...'")
            return None

        return shape_data

    except (AttributeError, TypeError) as e:
        if logger:
            logger.error(f"Shape data validation failed: {e}")
        return None


def _validate_required_properties(shape_data: Dict[str, Any], logger: Logger = None) -> bool:
    """Validate that all required properties exist."""
    missing_props = [
        prop for prop in PowerPointConfig.REQUIRED_PROPERTIES if prop not in shape_data
    ]

    if missing_props:
        if logger:
            logger.error(f"Missing required properties: {', '.join(missing_props)}, shape_data keys: {list(shape_data.keys())}, text='{str(shape_data.get('text', ''))[:30]}...'")
        return False
    
    return True


def _validate_and_convert_shape_type(shape_data: Dict[str, Any], logger: Logger = None) -> bool:
    """Validate and convert shape_type from string to enum or special types."""
    shape_type = shape_data.get("shape_type")
    
    if not isinstance(shape_type, str):
        if logger:
            logger.error("shape_type must be a string")
        return False

    shape_type_upper = shape_type.upper()

    cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
    if shape_type_upper == cfg["SHAPE_TYPE_TEXTBOX"]:
        shape_data["shape_type"] = cfg["SHAPE_TYPE_TEXTBOX"]
        return True

    if shape_type_upper == cfg["SHAPE_TYPE_AI_ICON"]:
        shape_data["shape_type"] = cfg["SHAPE_TYPE_AI_ICON"]
        return True

    try:
        shape_data["shape_type"] = getattr(MSO_AUTO_SHAPE_TYPE, shape_type_upper)
        return True
    except AttributeError:
        if logger:
            logger.error(f"Invalid shape_type: {shape_type}")
        return False


def _validate_numeric_properties(shape_data: Dict[str, Any], logger: Logger = None) -> bool:
    """Validate numeric properties efficiently."""
    # Validate positive numeric properties
    for prop in PowerPointConfig.NUMERIC_POSITIVE_PROPERTIES:
        if prop in shape_data:
            value = shape_data[prop]
            if not isinstance(value, (int, float)) or value < 0:
                if logger:
                    logger.error(f"Invalid {prop}: must be positive number, got {value}")
                return False

    # Validate coordinates
    for coord in ["x", "y"]:
        if coord in shape_data:
            value = shape_data[coord]
            if not isinstance(value, (int, float)):
                if logger:
                    logger.error(f"Invalid {coord} coordinate: must be number, got {value}")
                return False
    
    return True


def _validate_shape_boundaries(shape_data: Dict[str, Any], logger: Logger = None) -> bool:
    """Validate shape fits within slide boundaries."""
    x, y = shape_data["x"], shape_data["y"]
    width, height = shape_data["width"], shape_data["height"]

    # Check horizontal boundaries
    if x + width <= 0 or x >= PowerPointConfig.SLIDE_WIDTH_INCHES:
        if logger:
            logger.error(f"Shape outside horizontal boundaries: x={x}, width={width}")
        return False

    # Check vertical boundaries
    if y + height <= 0 or y >= PowerPointConfig.SLIDE_HEIGHT_INCHES:
        if logger:
            logger.error(f"Shape outside vertical boundaries: y={y}, height={height}")
        return False
    
    return True


def _validate_optional_properties(shape_data: Dict[str, Any], logger: Logger = None, max_font_size: Optional[int] = None) -> bool:
    """Validate optional properties like font_size and fill_color."""
    # Use custom max_font_size if provided, otherwise use default
    effective_max_font_size = max_font_size if max_font_size is not None else PowerPointConfig.MAX_FONT_SIZE
    
    # Validate font_size
    if "font_size" in shape_data and shape_data["font_size"] is not None:
        font_size = shape_data["font_size"]
        if not isinstance(font_size, (int, float)):
            if logger:
                logger.error(f"Invalid font_size: must be number, got {type(font_size).__name__}: {font_size}")
            return False
        if (
            font_size < PowerPointConfig.MIN_FONT_SIZE
            or font_size > effective_max_font_size
        ):
            if logger:
                logger.error(f"font_size {font_size} must be between {PowerPointConfig.MIN_FONT_SIZE} and {effective_max_font_size}")
            return False
    
    # Validate fill_color format
    if "fill_color" in shape_data and shape_data["fill_color"] is not None:
        fill_color = shape_data["fill_color"]
        if (
            not isinstance(fill_color, (list, tuple))
            or len(fill_color) != PowerPointConfig.RGB_COMPONENTS_COUNT
        ):
            if logger:
                logger.error(f"Invalid fill_color: must be RGB tuple/list with 3 values")
            return False

        # Convert to tuple if it's a list
        if isinstance(fill_color, list):
            shape_data["fill_color"] = tuple(fill_color)

        for color_value in fill_color:
            if not isinstance(color_value, int) or not (
                PowerPointConfig.RGB_MIN_VALUE
                <= color_value
                <= PowerPointConfig.RGB_MAX_VALUE
            ):
                if logger:
                    logger.error(f"Invalid fill_color: RGB values must be 0-255")
                return False

    # Validate text if present
    if "text" in shape_data:
        text = shape_data["text"]
        if text is not None and not isinstance(text, str):
            if logger:
                logger.error(f"Invalid text: must be string")
            return False
    
    # Validate font_color if present
    if "font_color" in shape_data and shape_data["font_color"] is not None:
        color_value = shape_data["font_color"]
        if not isinstance(color_value, (list, tuple)) or len(color_value) != 3:
            if logger:
                logger.error(f"Invalid font_color: must be RGB tuple/list with 3 values, got {type(color_value).__name__}: {color_value}")
            return False

        for val in color_value:
            if not isinstance(val, int) or not (0 <= val <= 255):
                if logger:
                    logger.error(f"Invalid font_color value: {val} (type={type(val).__name__}), full color={color_value}")
                if logger:
                    logger.error(f"Invalid font_color: RGB values must be 0-255")
                return False
    
    return True


# Property handlers - define how to handle each property
def apply_property_handlers():
    """Return dictionary of handler functions for each property"""

    def handle_text(shape, value, shape_data=None, slide=None):
        """
        Handle text property - DEPRECATED.

        Text is now handled directly in apply_text_to_shape().
        This function is kept for backward compatibility only.
        """
        return None
    
    def handle_fill_color(shape, value):
        if hasattr(shape, "fill") and value:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*value)

    def handle_transparency(shape, value):
        def _handle_transparency(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element
        if value:
            ts = shape.fill._xPr.solidFill
            sF = ts.get_or_change_to_srgbClr()
            _handle_transparency(sF, 'a:alpha', val=str(value))

    def handle_line_color(shape, value):
        if hasattr(shape, "line") and value:
            shape.line.fill.solid()
            shape.line.color.rgb = RGBColor(*value)

    def handle_font_bold(shape, value):
        if hasattr(shape, "text_frame") and value is not None:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = bool(value)

    def handle_text_color(shape, value):
        if hasattr(shape, "text_frame") and value:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*value)

    def handle_font_size(shape, value):
        if hasattr(shape, "text_frame") and value:
            text_frame = shape.text_frame
            # Apply font size to all paragraphs and runs
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(value)

    def handle_font_name(shape, value):
        if hasattr(shape, "text_frame") and value:
            text_frame = shape.text_frame
            # Apply font name to all paragraphs and runs
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = str(value)

    def handle_font_style(shape, value):
        if hasattr(shape, "text_frame") and value:
            text_frame = shape.text_frame
            # Apply font style to all paragraphs and runs
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if "bold" in value:
                        run.font.bold = True
                    if "italic" in value:
                        run.font.italic = True

    def handle_font_color(shape, value):
        if hasattr(shape, "text_frame") and value:
            text_frame = shape.text_frame
            # Apply font color to all paragraphs and runs
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*value)

    def handle_text_align(shape, value):
        if hasattr(shape, "text_frame") and value:
            text_frame = shape.text_frame
            # Apply text alignment to all paragraphs
            for paragraph in text_frame.paragraphs:
                if value == "center":
                    paragraph.alignment = PP_ALIGN.CENTER
                elif value == "left":
                    paragraph.alignment = PP_ALIGN.LEFT
                elif value == "right":
                    paragraph.alignment = PP_ALIGN.RIGHT

    def handle_vertical_align(shape, value):
        if hasattr(shape, "text_frame") and value:
            text_frame = shape.text_frame
            if value == "middle" or value == "center":
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            elif value == "top":
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            elif value == "bottom":
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM

    def handle_border_color(shape, value):
        if hasattr(shape, "line") and value:
            shape.line.color.rgb = RGBColor(*value)

    def handle_border_width(shape, value):
        if not hasattr(shape, "line"):
            return
        if value is None or value <= 0:
            shape.line.width = Pt(0)
            shape.line.fill.background()
        else:
            shape.line.width = Pt(value)

    def handle_corner_radius(shape, value):
        if (
            hasattr(shape, "auto_shape_type")
            and hasattr(shape, "adjustments")
            and value
        ):
            if shape.auto_shape_type != MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
                return
            # Normalize value to 0.0-1.0 ratio
            ratio = None
            if isinstance(value, str) and value.strip().endswith("%"):
                ratio = float(value.strip().rstrip("%")) / 100.0
            elif isinstance(value, (int, float)):
                ratio = float(value) / 100.0 if value > 1 else float(value)
            else:
                return
            # Clamp
            ratio = max(0.0, min(1.0, ratio))
            shape.adjustments[0] = ratio
            
    def handle_rotation(shape, value):
        """Handle rotation (in degrees, clockwise)."""
        try:
            if value is not None:
                shape.rotation = float(value)
        except Exception:
            pass

    # Return mapping from property name -> handler function
    return {
        "text": handle_text,
        "font_size": handle_font_size,
        "font_name": handle_font_name,
        "font_style": handle_font_style,
        "font_color": handle_font_color,
        "text_align": handle_text_align,
        "vertical_align": handle_vertical_align,
        "border_color": handle_border_color,
        "border_width": handle_border_width,
        "corner_radius": handle_corner_radius,
        "fill_color": handle_fill_color,
        "transparency": handle_transparency,
        "line_color": handle_line_color,
        "font_bold": handle_font_bold,
        "text_color": handle_text_color,
        "rotation": handle_rotation,
        # TODO: add other properties
    }


def get_all_supported_properties() -> List[str]:
    """Return all supported shape properties."""
    all_props = SHAPE_PROPERTIES["required"].copy()
    for category in SHAPE_PROPERTIES:
        if category != "required" and isinstance(SHAPE_PROPERTIES[category], dict):
            all_props.extend(SHAPE_PROPERTIES[category].keys())
    return all_props


def add_shapes_to_slide(
    shapes_data: List[Dict[str, Any]],
    prs: Presentation,
    slide_index: Optional[int] = None,
    output_file: Optional[str] = None,
    logger: Logger = None,
    max_font_size: Optional[int] = None,
) -> Presentation:
    """
    Create PowerPoint file with multiple shapes in a single slide.

    Uses dynamic property mapping to support flexible shape configuration.
    The PowerPoint file will be created with 16:9 aspect ratio.

    Args:
        shapes_data: List of dictionaries, each containing shape properties
        prs: Presentation object
        slide_index: Index of slide to add shapes to
        output_file: Output file name for the PowerPoint file
        logger: Logger instance for error logging
        max_font_size: Optional custom max font size (for PDF2PPTX: 72)

    Returns:
        Presentation object

    Supported properties can be found via get_all_supported_properties()
    """

    # Use a set for efficient lookups
    required_fields = set(PowerPointConfig.REQUIRED_PROPERTIES)

    # Filter out empty optional properties from each shape
    new_shapes_data = []
    for shape in shapes_data:
        filtered_shape = {}
        for key, value in shape.items():
            # Check if value should be included
            if key in required_fields:
                filtered_shape[key] = value
            elif value is not None:
                # Handle array-like values (list, tuple, numpy array)
                # Avoid comparing arrays with strings to prevent "ambiguous truth value" error
                if isinstance(value, (list, tuple)):
                    if len(value) > 0:
                        filtered_shape[key] = value
                else:
                    # Check if it's a numpy array
                    try:
                        import numpy as np
                        if isinstance(value, np.ndarray):
                            if value.size > 0:
                                filtered_shape[key] = value
                        elif isinstance(value, str):
                            # Only compare strings with empty string
                            if value != "":
                                filtered_shape[key] = value
                        else:
                            # For other types (int, float, bool, etc.), include if not None
                            filtered_shape[key] = value
                    except (ImportError, TypeError):
                        # If numpy is not available, check type before comparison
                        if isinstance(value, str):
                            if value != "":
                                filtered_shape[key] = value
                        else:
                            # For non-string types, include if not None
                            filtered_shape[key] = value
        new_shapes_data.append(filtered_shape)
    shapes_data = new_shapes_data

    # Set slide size to 16:9 aspect ratio
    prs.slide_width = Inches(PowerPointConfig.SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(PowerPointConfig.SLIDE_HEIGHT_INCHES)

    if slide_index is not None:
        slide = prs.slides[slide_index]
    else:
        slide = prs.slides.add_slide(
            prs.slide_layouts[PowerPointConfig.BLANK_SLIDE_LAYOUT_INDEX]
        )  # blank slide

    # Get property handlers
    property_handlers = apply_property_handlers()

    for shape_data in shapes_data:
        try:
            # Validate shape data format and values, get processed data
            processed_shape_data = validate_shape_data(shape_data.copy(), logger, max_font_size=max_font_size)
            
            # Skip if validation failed
            if processed_shape_data is None:
                if logger:
                    logger.warning(f"Skipping invalid shape: {shape_data.get('shape_type', 'unknown')}")
                continue

            # Get basic properties
            shape_type = processed_shape_data["shape_type"]
            original_width = processed_shape_data["width"]
            original_height = processed_shape_data["height"]
            x = processed_shape_data["x"]
            y = processed_shape_data["y"]

            # Get shape type name for text area ratio calculation
            cfg = PowerPointConfig.OBJECT_FUNCTION_CONFIG
            shape_type_name = shape_data.get("shape_type", cfg["DEFAULT_SHAPE_TYPE"])

            # NEW APPROACH: Use single shape or textbox with direct text_frame (full bounding box)
            if shape_type == cfg["SHAPE_TYPE_TEXTBOX"]:
                shape = slide.shapes.add_textbox(
                    Inches(x),
                    Inches(y),
                    Inches(original_width),
                    Inches(original_height),
                )
            # TODO: アイコンはBlobからではなく、nano banana proで生成した画像を用いる。(hiratani)
            elif shape_type == cfg["SHAPE_TYPE_AI_ICON"]:
                # Handle AI-generated icon insertion
                try:
                    from io import BytesIO

                    # from services.azure_service import AzureBlobService

                    # icon_name = shape_data.get("icon_name", cfg["DEFAULT_ICON_NAME"])
                    # issues_cfg = PowerPointConfig.ISSUES_DETAILS_CONFIG
                    # blob_rel_path = f"{issues_cfg[cfg['ICON_BLOB_PATH_KEY']]}/{icon_name}"

                    # azure_blob_service = AzureBlobService()
                    # blob_client = azure_blob_service.container_client.get_blob_client(
                    #     blob_rel_path
                    # )
                    # data = blob_client.download_blob().readall()
                    # bio = BytesIO(data)
                    # bio.seek(0)

                    # Insert picture instead of shape
                    shape = slide.shapes.add_picture(
                        shape_data.get("image_path"),
                        Inches(x),
                        Inches(y),
                        width=Inches(original_width),
                        height=Inches(original_height),
                    )

                    # 図形の陰影効果を消す（テーマ継承を切る）
                    if hasattr(shape, "shadow"):
                        try:
                            shape.shadow.inherit = False
                        except Exception:
                            pass
                    # Skip further processing since it's an image
                    continue

                except Exception as e:
                    # Fallback to text shape if icon loading fails
                    shape = slide.shapes.add_shape(
                        getattr(MSO_AUTO_SHAPE_TYPE, cfg["SHAPE_TYPE_ROUNDED_RECTANGLE"]),
                        Inches(x),
                        Inches(y),
                        Inches(original_width),
                        Inches(original_height),
                    )
                    # Add fallback text
                    text_frame = shape.text_frame
                    text_frame.text = cfg["FALLBACK_ICON_TEXT"]
                    continue
    
            else:
                shape = slide.shapes.add_shape(
                    shape_type,
                    Inches(x),
                    Inches(y),
                    Inches(original_width),
                    Inches(original_height),
                )

            # 図形の陰影効果を消す（テーマ継承を切る）
            if hasattr(shape, "shadow"):
                try:
                    shape.shadow.inherit = False
                except Exception:
                    pass

            line = shape.line
            line.fill.background()   # 枠線を完全に消す
            line.width = 0           # 念のため幅も0にする（環境差対策）

            # Set background color if provided (otherwise use PowerPoint default)
            if (
                "fill_color" in processed_shape_data
                and processed_shape_data["fill_color"] is not None
            ):
                fill_color = processed_shape_data["fill_color"]
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*fill_color)
                if (
                    "transparency" in processed_shape_data
                    and processed_shape_data["transparency"] is not None
                ):
                    transparency = int((1.0-processed_shape_data["transparency"])*100000)
                    processed_shape_data["transparency"] = transparency
            else:
                fill = shape.fill
                fill.background()

            # Set line color (tukawa)
            if (
                "line_color" in processed_shape_data
                and processed_shape_data["line_color"] is not None
            ):
                line_color = processed_shape_data["line_color"]
                shape.line.fill.solid()
                shape.line.color.rgb = RGBColor(*line_color)
            else:
                shape.line.fill.background()

            if "border_width" not in processed_shape_data:
                processed_shape_data["border_width"] = (
                    PowerPointConfig.DEFAULT_BORDER_WIDTH
                )

            # Handle text if provided - using full bounding box approach
            if "text" in processed_shape_data and processed_shape_data["text"]:
                text_value = processed_shape_data["text"]
                initial_font_size = processed_shape_data.get(
                    "font_size", PowerPointConfig.DEFAULT_FONT_SIZE
                )
                if logger:
                    logger.info(f"[add_shapes] Adding text shape: text='{text_value[:30]}...', x={x:.2f}, y={y:.2f}, w={original_width:.2f}, h={original_height:.2f}, font_size={initial_font_size}, font_color={processed_shape_data.get('font_color')}")

                # For TEXTBOX, automatically disable auto-resize unless explicitly enabled
                if shape_type != cfg["SHAPE_TYPE_TEXTBOX"]:
                    if "font_size" in processed_shape_data and processed_shape_data["font_size"]:
                        # JSONで指定されているフォントサイズをそのまま使う
                        font_size = processed_shape_data["font_size"]
                        optimal_result = {
                            "font_size": font_size,
                            "text_truncated": False,
                            "truncated_text": text_value,
                        }
                    else:
                        # 自動調整する
                        optimal_result = calculate_optimal_text_and_shape(
                            text=text_value,
                            original_width_inches=original_width,
                            original_height_inches=original_height,
                            shape_type_name=shape_type_name,
                            initial_font_size=initial_font_size,
                        )
                   
                    processed_shape_data["font_size"] = optimal_result["font_size"]

                    # Use truncated text if text was truncated to fit
                    if optimal_result.get("text_truncated", False):
                        processed_shape_data["text"] = optimal_result["truncated_text"]
                        text_value = optimal_result["truncated_text"]

                # Apply text to shape with 0.1 inch padding
                apply_text_to_shape(shape, text_value, processed_shape_data, logger)
                    
                if "rotation" in processed_shape_data:
                    try:
                        shape.rotation = float(processed_shape_data["rotation"])
                    except Exception:
                        pass
    
            # Apply all other properties
            for prop_name, prop_value in processed_shape_data.items():
                # Skip already processed properties
                if prop_name in SHAPE_PROPERTIES["required"] or prop_name == "text":
                    continue
                elif prop_name == "arrow":
                    continue
                # Find and run corresponding handler
                if prop_name in property_handlers:
                    try:
                        property_handlers[prop_name](shape, prop_value)
                    except Exception as e:
                        if logger:
                            logger.warning(f"Cannot apply property '{prop_name}': {e}")
                else:
                    if logger:
                        logger.debug(f"No handler found for property: '{prop_name}'")

        except Exception as e:
            if logger:
                logger.warning(f"Cannot add shape: '{e}'")
            continue

    if output_file:
        prs.save(output_file)

    return prs
