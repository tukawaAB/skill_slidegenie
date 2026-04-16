"""Slide generation pipeline — orchestrates the 4-step process.

Ported from: ppt-addin/backend/services/make_pptx/make_pptx_service.py
"""
import os
from pathlib import Path

from PIL import Image

from slidegenie.auth import get_genai_client
from slidegenie.gemini_client import GEMINIClient
from slidegenie.image_gen.builder import prompt_to_image
from slidegenie.json_gen.builder import image_to_json
from slidegenie.json_gen.postprocess import json_postprocess
from slidegenie.slide_gen.builder import json_to_pptx, multi_json_to_pptx
from slidegenie.utils.common import load_prompt, is_english
from slidegenie.utils.constants import PowerPointConfig
from slidegenie.utils.logger import get_logger

_gemini_client: GEMINIClient | None = None


def _get_gemini_client(logger=None) -> GEMINIClient:
    """Get or create a singleton GEMINIClient."""
    global _gemini_client
    if _gemini_client is None:
        genai_client = get_genai_client()
        _gemini_client = GEMINIClient(client=genai_client, logger=logger)
    return _gemini_client


def select_make_type(
    prompt: str,
    gemini_client: GEMINIClient,
    make_types: list[str] | None = None,
    image: Image.Image | None = None,
) -> str:
    """Select optimal slide type using Gemini AI."""
    if make_types is None:
        make_types = ["graphic", "flow", "matrix"]

    prompt_text = load_prompt(
        "select_maketype",
        user_prompt=prompt,
        make_types=make_types,
    )
    result = gemini_client.generate_chat_completion(
        prompt=prompt_text,
        image=image,
    )

    # Extract make_type from response
    result = result.strip().lower()
    for mt in make_types:
        if mt in result:
            return mt

    # Default fallback
    return make_types[0]


def generate_slide(
    prompt: str,
    output_path: str = "output.pptx",
    mode: str = "editable",
    make_type: str | None = None,
    make_types: list[str] | None = None,
    image: Image.Image | None = None,
    tone_manner: str | None = None,
) -> dict:
    """Generate a single slide.

    Args:
        prompt: Content description for the slide
        output_path: Path for the output file
        mode: "image", "editable", or "both"
        make_type: Explicit slide type (graphic/flow/matrix), auto-select if None
        make_types: Available types for auto-selection
        image: Optional input image for reference
        tone_manner: Optional tone & manner text

    Returns:
        dict with keys: image_path, pptx_path (depending on mode)
    """
    logger = get_logger("pipeline")
    gemini_client = _get_gemini_client(logger)
    result = {}

    # Step 1: Language detection
    english_frag = is_english(prompt)
    lang = "en" if english_frag else "ja"
    logger.info(f"Language detected: {lang}")

    # Step 2: Make type selection
    if make_type is None:
        make_type = select_make_type(
            prompt=prompt,
            gemini_client=gemini_client,
            make_types=make_types,
            image=image,
        )
    logger.info(f"Make type: {make_type}")

    # Step 3: Image generation
    logger.info("Generating slide image...")
    slide_image = prompt_to_image(
        make_type=make_type,
        prompt=prompt,
        logger=logger,
        gemini_client=gemini_client,
        english_frag=english_frag,
        image=image,
        tone_manner=tone_manner,
    )
    logger.info("Image generation complete")

    # Save image if requested
    if mode in ("image", "both"):
        image_path = str(Path(output_path).with_suffix(".png"))
        slide_image.save(image_path, format="PNG")
        logger.info(f"Image saved: {image_path}")
        result["image_path"] = image_path

    # Step 4: OCR + PPTX if requested
    if mode in ("editable", "both"):
        # OCR: image → JSON
        logger.info("Running OCR...")
        json_data = image_to_json(
            make_type=make_type,
            logger=logger,
            gemini_client=gemini_client,
            prompt=prompt,
            src_image=slide_image,
        )

        # Post-process JSON
        json_data = json_postprocess(json_data)

        # JSON → PPTX
        pptx_path = str(Path(output_path).with_suffix(".pptx"))
        json_to_pptx(
            make_type=make_type,
            logger=logger,
            json_data=json_data,
            src_image=slide_image,
            output_path=pptx_path,
        )
        result["pptx_path"] = pptx_path

    # If image-only mode, create a PPTX with just the image
    if mode == "image":
        pptx_path = str(Path(output_path).with_suffix(".pptx"))
        _create_image_only_pptx(slide_image, pptx_path, logger)
        result["pptx_path"] = pptx_path

    return result


def convert_image(
    input_path: str,
    output_path: str = "output.pptx",
) -> dict:
    """Convert an existing image to editable PPTX.

    Args:
        input_path: Path to the input image
        output_path: Path for the output PPTX

    Returns:
        dict with key: pptx_path
    """
    logger = get_logger("pipeline")
    gemini_client = _get_gemini_client(logger)

    # Load image
    src_image = Image.open(input_path)
    logger.info(f"Loaded image: {input_path} ({src_image.size})")

    # OCR: image → JSON
    logger.info("Running OCR...")
    json_data = image_to_json(
        make_type="graphic",
        logger=logger,
        gemini_client=gemini_client,
        src_image=src_image,
    )

    # Post-process
    json_data = json_postprocess(json_data)

    # JSON → PPTX
    pptx_path = str(Path(output_path).with_suffix(".pptx"))
    json_to_pptx(
        make_type="graphic",
        logger=logger,
        json_data=json_data,
        src_image=src_image,
        output_path=pptx_path,
    )

    return {"pptx_path": pptx_path}


def _generate_slide_image_only(
    prompt: str,
    make_type: str | None = None,
    make_types: list[str] | None = None,
    image: Image.Image | None = None,
    tone_manner: str | None = None,
    save_image_path: str | None = None,
) -> dict:
    """Generate image only (no OCR). For image-mode batch processing.

    Returns:
        dict with keys: src_image, make_type
    """
    logger = get_logger("pipeline")
    gemini_client = _get_gemini_client(logger)

    english_frag = is_english(prompt)

    if make_type is None:
        make_type = select_make_type(
            prompt=prompt, gemini_client=gemini_client,
            make_types=make_types, image=image,
        )
    logger.info(f"Make type: {make_type}")

    slide_image = prompt_to_image(
        make_type=make_type, prompt=prompt, logger=logger,
        gemini_client=gemini_client, english_frag=english_frag,
        image=image, tone_manner=tone_manner,
    )

    if save_image_path:
        slide_image.save(save_image_path, format="PNG")

    return {"src_image": slide_image, "make_type": make_type}


def _generate_slide_json(
    prompt: str,
    make_type: str | None = None,
    make_types: list[str] | None = None,
    image: Image.Image | None = None,
    tone_manner: str | None = None,
    save_image_path: str | None = None,
) -> dict:
    """Generate image + OCR → return JSON data and image (no PPTX creation).

    Returns:
        dict with keys: json_data, src_image, make_type
    """
    logger = get_logger("pipeline")
    gemini_client = _get_gemini_client(logger)

    english_frag = is_english(prompt)

    if make_type is None:
        make_type = select_make_type(
            prompt=prompt, gemini_client=gemini_client,
            make_types=make_types, image=image,
        )
    logger.info(f"Make type: {make_type}")

    # Image generation
    slide_image = prompt_to_image(
        make_type=make_type, prompt=prompt, logger=logger,
        gemini_client=gemini_client, english_frag=english_frag,
        image=image, tone_manner=tone_manner,
    )

    if save_image_path:
        slide_image.save(save_image_path, format="PNG")

    # OCR → JSON
    json_data = image_to_json(
        make_type=make_type, logger=logger,
        gemini_client=gemini_client, prompt=prompt, src_image=slide_image,
    )
    json_data = json_postprocess(json_data)

    return {"json_data": json_data, "src_image": slide_image, "make_type": make_type}


def batch_generate(
    specs: list[dict],
    output_path: str,
    mode: str = "editable",
    max_workers: int = 4,
    image_dir: str | None = None,
) -> str:
    """Generate multiple slides in parallel and build a single PPTX.

    Phase 1: Parallel image generation + OCR → JSON for each slide
    Phase 2: Sequential JSON → PPTX (single file, all slides)

    Args:
        specs: List of dicts with keys: prompt, make_type (optional)
        output_path: Path for the output PPTX
        mode: Output mode (image/editable/both)
        max_workers: Maximum parallel workers for phase 1
        image_dir: Directory to save slide images (optional)

    Returns:
        Path to the generated PPTX
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed

    logger = get_logger("pipeline")

    # Pre-initialize the Gemini client (singleton, thread-safe)
    _get_gemini_client(logger)

    if image_dir:
        os.makedirs(image_dir, exist_ok=True)

    # Phase 1: Parallel generation
    slide_results = [None] * len(specs)

    if mode == "image":
        # Image-only mode: skip OCR, just generate images in parallel
        logger.info(f"Phase 1: Generating {len(specs)} slide images in parallel (workers={max_workers})...")

        def _gen_image(idx, spec):
            logger.info(f"[{idx + 1}/{len(specs)}] Generating image...")
            img_path = None
            if image_dir:
                img_path = os.path.join(image_dir, f"slide_{idx + 1}.png")
            result = _generate_slide_image_only(
                prompt=spec["prompt"],
                make_type=spec.get("make_type"),
                save_image_path=img_path,
            )
            logger.info(f"[{idx + 1}/{len(specs)}] Image ready")
            return idx, result

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(_gen_image, i, s): i for i, s in enumerate(specs)}
            for future in as_completed(futures):
                idx, result = future.result()
                slide_results[idx] = result
    else:
        # Editable/both mode: need OCR
        logger.info(f"Phase 1: Generating {len(specs)} slides in parallel (workers={max_workers})...")

        def _gen_json(idx, spec):
            logger.info(f"[{idx + 1}/{len(specs)}] Generating JSON...")
            img_path = None
            if image_dir:
                img_path = os.path.join(image_dir, f"slide_{idx + 1}.png")
            result = _generate_slide_json(
                prompt=spec["prompt"],
                make_type=spec.get("make_type"),
                save_image_path=img_path,
            )
            logger.info(f"[{idx + 1}/{len(specs)}] JSON ready: {result['json_data'].get('title', '')[:40]}")
            return idx, result

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(_gen_json, i, s): i for i, s in enumerate(specs)}
            for future in as_completed(futures):
                idx, result = future.result()
                slide_results[idx] = result

    # Phase 2: Build single PPTX with all slides
    logger.info(f"Phase 2: Building PPTX with {len(specs)} slides...")

    if mode in ("editable", "both"):
        slides_data = []
        for r in slide_results:
            slides_data.append({
                "json_data": r["json_data"],
                "src_image": r["src_image"],
            })
        multi_json_to_pptx(
            slides_data=slides_data,
            logger=logger,
            output_path=output_path,
        )
    elif mode == "image":
        # Image-only: paste each image as a full slide
        from pptx import Presentation as Prs
        from pptx.util import Emu
        from io import BytesIO

        prs = Prs()
        prs.slide_width = Emu(int(PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.EMU_PER_INCH))
        prs.slide_height = Emu(int(PowerPointConfig.SLIDE_HEIGHT_INCHES * PowerPointConfig.EMU_PER_INCH))

        for r in slide_results:
            slide = prs.slides.add_slide(prs.slide_layouts[min(len(prs.slide_layouts) - 1, 6)])
            img_bytes = BytesIO()
            r["src_image"].save(img_bytes, format="PNG")
            img_bytes.seek(0)
            slide.shapes.add_picture(
                img_bytes, 0, 0,
                width=Emu(int(PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.EMU_PER_INCH)),
                height=Emu(int(PowerPointConfig.SLIDE_HEIGHT_INCHES * PowerPointConfig.EMU_PER_INCH)),
            )

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)

    logger.info(f"Done: {output_path}")
    return output_path


def _create_image_only_pptx(
    image: Image.Image, output_path: str, logger
) -> str:
    """Create a PPTX with just the image pasted as full slide."""
    from pptx import Presentation as Prs
    from pptx.util import Inches, Emu

    prs = Prs()
    prs.slide_width = Emu(int(PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.EMU_PER_INCH))
    prs.slide_height = Emu(int(PowerPointConfig.SLIDE_HEIGHT_INCHES * PowerPointConfig.EMU_PER_INCH))

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    from io import BytesIO
    img_bytes = BytesIO()
    image.save(img_bytes, format="PNG")
    img_bytes.seek(0)

    slide.shapes.add_picture(
        img_bytes, 0, 0,
        width=Emu(int(PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.EMU_PER_INCH)),
        height=Emu(int(PowerPointConfig.SLIDE_HEIGHT_INCHES * PowerPointConfig.EMU_PER_INCH)),
    )

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    logger.info(f"Image-only PPTX saved: {output_path}")
    return output_path


def merge_pptx(input_paths: list[str], output_path: str) -> str:
    """Merge multiple PPTX files into one, correctly copying embedded images.

    Args:
        input_paths: List of paths to source PPTX files
        output_path: Path for the merged output PPTX

    Returns:
        Path to the merged PPTX
    """
    from pptx import Presentation as Prs
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from io import BytesIO

    logger = get_logger("pipeline")

    prs = Prs()
    prs.slide_width = Emu(int(PowerPointConfig.SLIDE_WIDTH_INCHES * PowerPointConfig.EMU_PER_INCH))
    prs.slide_height = Emu(int(PowerPointConfig.SLIDE_HEIGHT_INCHES * PowerPointConfig.EMU_PER_INCH))

    blank_layout = prs.slide_layouts[min(len(prs.slide_layouts) - 1, 6)]

    for input_path in input_paths:
        logger.info(f"Merging: {input_path}")
        src = Prs(input_path)
        for slide in src.slides:
            new_slide = prs.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract image blob and re-embed in new presentation
                    img_bytes = BytesIO(shape.image.blob)
                    new_slide.shapes.add_picture(
                        img_bytes,
                        shape.left, shape.top,
                        shape.width, shape.height,
                    )
                elif shape.has_text_frame:
                    # Copy text shapes
                    from pptx.util import Inches
                    txbox = new_slide.shapes.add_textbox(
                        shape.left, shape.top,
                        shape.width, shape.height,
                    )
                    for i, para in enumerate(shape.text_frame.paragraphs):
                        if i > 0:
                            txbox.text_frame.add_paragraph()
                        p = txbox.text_frame.paragraphs[i]
                        p.alignment = para.alignment
                        for run in para.runs:
                            r = p.add_run()
                            r.text = run.text
                            if run.font.size:
                                r.font.size = run.font.size
                            if run.font.bold is not None:
                                r.font.bold = run.font.bold
                            if run.font.color and run.font.color.rgb:
                                r.font.color.rgb = run.font.color.rgb

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    logger.info(f"Merged {len(input_paths)} files → {output_path}")
    return output_path
